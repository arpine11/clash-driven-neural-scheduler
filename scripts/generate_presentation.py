"""
generate_presentation.py
------------------------
Generates the 10-slide capstone defense presentation as a .pptx file.
Output: outputs/reports/Tadevosyan_AUA_CS_Capstone_2026_Presentation.pptx

Run:
    .venv/bin/python3 scripts/generate_presentation.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT      = Path(__file__).resolve().parent.parent
FIGS      = ROOT / "figs"
DATA_VIS  = ROOT / "data_visuals"
SENS      = ROOT / "sensitivity_figs"
MIN_TS    = ROOT / "outputs" / "min_timeslot_experiment"
EXT_PLOTS = MIN_TS / "extended_dynamics" / "plots"
ROLL      = MIN_TS / "clash_dynamics" / "rolling_plots"
OUT_DIR   = ROOT / "outputs" / "reports"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH  = OUT_DIR / "Tadevosyan_AUA_CS_Capstone_2026_Presentation.pptx"

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
NAVY      = RGBColor(0x1A, 0x2A, 0x4A)   # deep navy — title bg, accents
GOLD      = RGBColor(0xC8, 0xA0, 0x32)   # gold — highlights, callout borders
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)   # near-white slide bg
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MID_GRAY  = RGBColor(0x55, 0x65, 0x7A)
DARK_TEXT = RGBColor(0x1C, 0x1C, 0x2E)
GREEN_OK  = RGBColor(0x2C, 0xA0, 0x2C)
RED_BAD   = RGBColor(0xD6, 0x27, 0x28)
BLUE_DYN  = RGBColor(0x4C, 0x72, 0xB0)
ORG_HYB   = RGBColor(0xDD, 0x84, 0x52)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def fill_shape(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def add_rect(slide, left, top, width, height, rgb: RGBColor, radius=0):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    fill_shape(shape, rgb)
    shape.line.fill.background()   # no border
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=DARK_TEXT,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True,
                 font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font_name
    return txb


def add_para(tf, text, font_size=16, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, space_before=6, italic=False,
             font_name="Calibri"):
    p   = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font_name
    return p


def add_image(slide, path: Path, left, top, width, height=None):
    if not path.exists():
        return None
    if height is None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    return slide.shapes.add_picture(str(path), left, top, width, height)


def header_bar(slide, title_text: str, subtitle: str = ""):
    """Dark navy top bar with slide title."""
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    # gold accent line beneath bar
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.04), GOLD)

    add_text_box(slide, title_text,
                 Inches(0.4), Inches(0.12), Inches(10), Inches(0.65),
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.72), Inches(10), Inches(0.38),
                     font_size=14, bold=False, color=RGBColor(0xCC, 0xD6, 0xE8),
                     align=PP_ALIGN.LEFT, italic=True)


def slide_bg(slide):
    """Light background rectangle covering full slide."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_BG)


def callout_box(slide, text, left, top, width, height,
                bg=RGBColor(0xE8, 0xF0, 0xFE), border=NAVY,
                font_size=15, bold=False, color=DARK_TEXT):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = border
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return box


def bullet_box(slide, items: list[tuple[str, str]],
               left, top, width, height,
               title_size=15, body_size=14):
    """
    items = list of (bullet_char, text)  e.g. ("•", "Some point")
    """
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for bchar, text in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bchar}  {text}"
        run.font.size  = Pt(body_size)
        run.font.color.rgb = DARK_TEXT
        run.font.name  = "Calibri"
    return txb


def footer(slide, text="Arpine Tadevosyan · AUA CS Capstone 2026"):
    add_rect(slide, 0, Inches(7.22), SLIDE_W, Inches(0.28), NAVY)
    add_text_box(slide, text,
                 Inches(0.3), Inches(7.23), Inches(12), Inches(0.24),
                 font_size=9, color=RGBColor(0xAA, 0xBB, 0xCC),
                 align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = blank_slide(prs)

    # Full navy background
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)

    # Gold horizontal accent band
    add_rect(s, 0, Inches(2.9), SLIDE_W, Inches(0.06), GOLD)

    # Main title
    txb = s.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(11.9), Inches(1.6))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "A Hybrid Propagating-Particle and Hopfield Network\nAlgorithm for Exam Timetabling"
    run.font.size  = Pt(32)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    run.font.name  = "Calibri"

    # Subtitle
    add_text_box(s, "CS Capstone Defense — American University of Armenia",
                 Inches(0.7), Inches(3.1), Inches(11.9), Inches(0.5),
                 font_size=17, color=RGBColor(0xCC, 0xD6, 0xE8),
                 align=PP_ALIGN.CENTER, italic=True)

    # Author block
    add_text_box(s, "Arpine Tadevosyan",
                 Inches(0.7), Inches(3.85), Inches(11.9), Inches(0.45),
                 font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_text_box(s, "Supervisors: Suren Khachatryan · Aleksandr Hayrapetyan",
                 Inches(0.7), Inches(4.4), Inches(11.9), Inches(0.4),
                 font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                 align=PP_ALIGN.CENTER)

    add_text_box(s, "May 2026",
                 Inches(0.7), Inches(4.9), Inches(11.9), Inches(0.4),
                 font_size=13, color=RGBColor(0x88, 0x99, 0xAA),
                 align=PP_ALIGN.CENTER, italic=True)

    # Keywords strip at bottom
    add_rect(s, 0, Inches(6.55), SLIDE_W, Inches(0.55), RGBColor(0x12, 0x1E, 0x38))
    add_text_box(s,
        "exam timetabling  ·  propagating particles  ·  Hopfield network  ·  "
        "discrete dynamical systems  ·  Toronto benchmark  ·  NP-complete",
        Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.45),
        font_size=10, color=RGBColor(0x77, 0x99, 0xBB),
        align=PP_ALIGN.CENTER, italic=True)


# ---------------------------------------------------------------------------
# Slide 2 — Problem Statement
# ---------------------------------------------------------------------------

def slide_problem(prs):
    s = blank_slide(prs)
    slide_bg(s)
    header_bar(s, "The Problem", "NP-complete — but can we understand it physically?")

    # Left column — bullets
    LEFT = Inches(0.4)
    TOP  = Inches(1.35)
    W    = Inches(5.8)

    items = [
        ("①", "Assign N exams to S timeslots"),
        ("②", "Constraint: two exams sharing a student cannot occupy the same slot"),
        ("③", "Model as a conflict graph → valid schedule = proper graph coloring"),
        ("④", "Graph coloring is NP-complete — no known efficient general solution"),
        ("⑤", "Universities run this every semester with hundreds of exams"),
    ]
    bullet_box(s, items, LEFT, TOP, W, Inches(3.5), body_size=15)

    # Insight callout
    callout_box(s,
        "Most tools: black-box heuristics.\n"
        "This project: treat scheduling as a physical dynamical system\n"
        "— and explain why it works.",
        LEFT, Inches(5.05), W, Inches(1.0),
        bg=RGBColor(0xE8, 0xF4, 0xE8), border=GREEN_OK,
        font_size=13, bold=False)

    # Right column — conflict graph illustration
    img_path = DATA_VIS / "car-f-92_clash_heatmap.png"
    add_image(s, img_path, Inches(6.5), Inches(1.35), Inches(6.5))

    add_text_box(s, "Conflict heatmap — car-f-92 (682 exams)",
                 Inches(6.5), Inches(6.5), Inches(6.5), Inches(0.4),
                 font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    footer(s)


# ---------------------------------------------------------------------------
# Slide 3 — Toronto Benchmark
# ---------------------------------------------------------------------------

def slide_dataset(prs):
    s = blank_slide(prs)
    slide_bg(s)
    header_bar(s, "The Toronto Benchmark",
               "12 real instances · Carter, Laporte & Lee (1996) · standard for 30 years")

    # Table data
    rows = [
        ("Instance",    "Exams", "Density", "Std. Slots", "Min. Slots", "Difficulty"),
        ("sta-f-83",    "139",   "14%",     "13",         "13",         "Easy"),
        ("ute-s-92",    "184",   "8%",      "10",         "10",         "Easy"),
        ("hec-s-92",    "81",    "42%",     "18",         "17",         "Medium"),
        ("ear-f-83",    "190",   "27%",     "24",         "22",         "Medium"),
        ("kfu-s-93",    "461",   "6%",      "20",         "19",         "Medium"),
        ("lse-f-91",    "381",   "6%",      "18",         "17",         "Medium"),
        ("rye-s-93",    "486",   "5%",      "23",         "21",         "Medium"),
        ("tre-s-92",    "261",   "6%",      "23",         "20",         "Hard"),
        ("yor-f-83",    "181",   "29%",     "21",         "19",         "Hard"),
        ("uta-s-92",    "622",   "8%",      "35",         "30",         "Hard"),
        ("car-f-92",    "682",   "13%",     "32",         "28",         "Hardest"),
        ("car-s-91",    "682",   "14%",     "35",         "28",         "Hardest"),
    ]

    col_widths = [Inches(1.3), Inches(0.75), Inches(0.85), Inches(1.05),
                  Inches(1.05), Inches(1.0)]
    row_height = Inches(0.37)

    table_left = Inches(0.35)
    table_top  = Inches(1.3)

    from pptx.util import Inches as I_
    tbl = s.shapes.add_table(
        len(rows), len(col_widths),
        table_left, table_top,
        sum(col_widths), row_height * len(rows)
    ).table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    diff_colors = {
        "Easy":    RGBColor(0xD4, 0xED, 0xDA),
        "Medium":  RGBColor(0xFF, 0xF3, 0xCD),
        "Hard":    RGBColor(0xF8, 0xD7, 0xDA),
        "Hardest": RGBColor(0xF5, 0xC6, 0xCB),
    }

    for ri, row_data in enumerate(rows):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.text = cell_text
            run.font.name = "Calibri"
            run.font.size = Pt(11) if ri > 0 else Pt(11)
            run.font.bold = (ri == 0)
            if ri == 0:
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                diff = row_data[5]
                run.font.color.rgb = DARK_TEXT
                cell.fill.solid()
                if ci == 5 and diff in diff_colors:
                    cell.fill.fore_color.rgb = diff_colors[diff]
                else:
                    cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else RGBColor(0xF2, 0xF5, 0xF9)

    # Right side: scatter plot
    img = DATA_VIS / "00_size_vs_density.png"
    add_image(s, img, Inches(6.6), Inches(1.28), Inches(6.4))
    add_text_box(s, "Instance size vs. conflict density",
                 Inches(6.6), Inches(6.45), Inches(6.4), Inches(0.35),
                 font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    footer(s)


# ---------------------------------------------------------------------------
# Slide 4 — The Physics Idea
# ---------------------------------------------------------------------------

def slide_physics(prs):
    s = blank_slide(prs)
    slide_bg(s)
    header_bar(s, "The Physics Idea",
               "The schedule is a point on a discrete torus — the algorithm moves it")

    # Left bullets
    items = [
        ("▸", "Current schedule = point on N-dimensional discrete torus"),
        ("▸", "N axes (one per exam) × S positions (timeslots) each"),
        ("▸", "E operator: shifts a clashing exam +1 slot, up to m times per step"),
        ("▸", "m = mode parameter — controls geometry of motion"),
    ]
    bullet_box(s, items, Inches(0.4), Inches(1.35), Inches(5.6), Inches(2.2),
               body_size=15)

    # Key insight box
    callout_box(s,
        "gcd(m, S) = 1  →  quasi-periodic  (ergodic: reaches any slot)\n"
        "gcd(m, S) > 1  →  resonant  (trapped in smaller orbit)",
        Inches(0.4), Inches(3.7), Inches(5.6), Inches(0.95),
        bg=RGBColor(0xE8, 0xF0, 0xFE), border=NAVY,
        font_size=14, bold=False)

    add_text_box(s,
        "Same idea as orbital resonance in mechanics —\napplied to a combinatorial scheduling space.",
        Inches(0.4), Inches(4.8), Inches(5.6), Inches(0.7),
        font_size=13, italic=True, color=MID_GRAY)

    # Figure: mode type comparison
    img = FIGS / "mode_type_comparison.png"
    add_image(s, img, Inches(6.3), Inches(1.28), Inches(6.7))
    add_text_box(s, "Performance: quasi-periodic vs. resonant modes across instances",
                 Inches(6.3), Inches(6.45), Inches(6.7), Inches(0.35),
                 font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    footer(s)


# ---------------------------------------------------------------------------
# Slide 5 — The Algorithm
# ---------------------------------------------------------------------------

def slide_algorithm(prs):
    s = blank_slide(prs)
    slide_bg(s)
    header_bar(s, "The Algorithm",
               "Dynamic-only walk + Hopfield network repair, with rollback safeguard")

    # Dynamic-only box
    dyn_box = s.shapes.add_shape(1, Inches(0.35), Inches(1.35), Inches(4.0), Inches(4.6))
    dyn_box.fill.solid()
    dyn_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFF)
    dyn_box.line.color.rgb = BLUE_DYN
    dyn_box.line.width = Pt(2)

    add_text_box(s, "Dynamic-Only", Inches(0.4), Inches(1.38),
                 Inches(3.9), Inches(0.4), font_size=15, bold=True, color=BLUE_DYN)

    dyn_items = [
        ("①", "Sort exams by conflict count (most conflicted first)"),
        ("②", "Apply E operator: shift clashing exam +1 slot"),
        ("③", "Repeat up to m times per step"),
        ("④", "Use best quasi-periodic mode per instance"),
        ("⑤", "No memory — pure physics walk"),
    ]
    bullet_box(s, dyn_items, Inches(0.45), Inches(1.85), Inches(3.8), Inches(2.8),
               body_size=13)

    # Hybrid box
    hyb_box = s.shapes.add_shape(1, Inches(4.65), Inches(1.35), Inches(4.7), Inches(4.6))
    hyb_box.fill.solid()
    hyb_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xE8)
    hyb_box.line.color.rgb = ORG_HYB
    hyb_box.line.width = Pt(2)

    add_text_box(s, "Hybrid (adds Hopfield)", Inches(4.7), Inches(1.38),
                 Inches(4.55), Inches(0.4), font_size=15, bold=True, color=ORG_HYB)

    hyb_items = [
        ("①", "Run E operator as above"),
        ("②", "Collect timetable snapshots during run"),
        ("③", "Train Hopfield network (Hebbian weights) on snapshots"),
        ("④", "When stuck: recall — pull clashing exams toward nearest stored pattern"),
        ("⑤", "Rollback safeguard: if Hopfield call worsens clashes → revert immediately"),
    ]
    bullet_box(s, hyb_items, Inches(4.7), Inches(1.85), Inches(4.55), Inches(3.5),
               body_size=13)

    # Hopfield effect figure
    img = FIGS / "hopfield_effect.png"
    add_image(s, img, Inches(9.6), Inches(1.28), Inches(3.45))
    add_text_box(s, "Effect of Hopfield repair step",
                 Inches(9.6), Inches(5.55), Inches(3.45), Inches(0.35),
                 font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    # Rollback callout at bottom
    callout_box(s,
        "Without rollback: Hopfield destabilises the descent.\n"
        "With rollback: it becomes a reliable escape mechanism.",
        Inches(0.35), Inches(6.08), Inches(8.85), Inches(0.72),
        bg=RGBColor(0xFF, 0xF3, 0xCD), border=GOLD,
        font_size=13, bold=False)

    footer(s)


# ---------------------------------------------------------------------------
# Slide 6 — Standard Timeslot Results
# ---------------------------------------------------------------------------

def slide_standard(prs):
    s = blank_slide(prs)
    slide_bg(s)
    header_bar(s, "Baseline Results: Standard Timeslots",
               "Under full schedules — dynamic-only alone is already a strong solver")

    # Big result callout
    callout_box(s,
        "10 / 12 instances solved to ZERO clashes — dynamic-only method",
        Inches(0.4), Inches(1.32), Inches(8.6), Inches(0.65),
        bg=RGBColor(0xD4, 0xED, 0xDA), border=GREEN_OK,
        font_size=17, bold=True, color=RGBColor(0x15, 0x55, 0x24))

    # Bullets
    items = [
        ("✓", "10 / 12 instances → zero clashes (dynamic-only)"),
        ("✗", "2 unsolved: car-f-92 & car-s-91 — largest & densest in benchmark"),
        ("✓", "Hybrid matches or improves all results"),
        ("✓", "Hopfield repair reduces residual clashes further on the 2 hard instances"),
        ("→", "Key insight: with enough timeslots, physics walk alone suffices"),
    ]
    bullet_box(s, items, Inches(0.4), Inches(2.15), Inches(5.6), Inches(2.8),
               body_size=15)

    callout_box(s,
        "The hybrid becomes most valuable when constraints are tightened.\n"
        "Which is exactly what we tested next →",
        Inches(0.4), Inches(5.2), Inches(5.6), Inches(0.8),
        bg=RGBColor(0xE8, 0xF0, 0xFE), border=NAVY,
        font_size=13)

    # Figure
    img = FIGS / "dyn_vs_hybrid.png"
    add_image(s, img, Inches(6.3), Inches(1.28), Inches(6.7))
    add_text_box(s, "Dynamic-only vs. Hybrid — final clash count across all 12 instances",
                 Inches(6.3), Inches(6.45), Inches(6.7), Inches(0.35),
                 font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    footer(s)


# ---------------------------------------------------------------------------
# Slide 7 — Minimum Timeslots
# ---------------------------------------------------------------------------

def slide_minimum(prs):
    s = blank_slide(prs)
    slide_bg(s)
    header_bar(s, "The Harder Challenge: Minimum Timeslots",
               "S reduced to theoretical minimum — right at the edge of feasibility")

    # Results table
    rows_data = [
        ("Instance",  "Δ (dyn − hyb)", "Winner"),
        ("car-s-91",  "+142",           "Hybrid ✓"),
        ("uta-s-92",  "+40",            "Hybrid ✓"),
        ("tre-s-92",  "+20",            "Hybrid ✓"),
        ("yor-f-83",  "+10",            "Hybrid ✓"),
        ("ear-f-83",  "+8",             "Hybrid ✓"),
        ("hec-s-92",  "0",              "Tie"),
        ("kfu-s-93",  "0",              "Tie"),
        ("rye-s-93",  "0",              "Tie"),
        ("sta-f-83",  "0",              "Tie (both 0)"),
        ("ute-s-92",  "0",              "Tie (both 0)"),
        ("lse-f-91",  "−2",             "Dynamic"),
        ("car-f-92",  "−4",             "Dynamic"),
    ]

    col_widths = [Inches(1.35), Inches(1.4), Inches(1.35)]
    row_h = Inches(0.35)

    tbl = s.shapes.add_table(
        len(rows_data), 3,
        Inches(0.35), Inches(1.32),
        sum(col_widths), row_h * len(rows_data)
    ).table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    winner_colors = {
        "Hybrid ✓": RGBColor(0xD4, 0xED, 0xDA),
        "Tie":       RGBColor(0xFF, 0xF3, 0xCD),
        "Tie (both 0)": RGBColor(0xD4, 0xED, 0xDA),
        "Dynamic":  RGBColor(0xF8, 0xD7, 0xDA),
    }

    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = cell_text
            run.font.name = "Calibri"
            run.font.size = Pt(11)
            run.font.bold = (ri == 0)
            if ri == 0:
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                winner = row_data[2]
                cell.fill.solid()
                if ci == 2 and winner in winner_colors:
                    cell.fill.fore_color.rgb = winner_colors[winner]
                else:
                    cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else RGBColor(0xF2, 0xF5, 0xF9)
                run.font.color.rgb = DARK_TEXT

    # Score summary
    score_items = [
        ("🟢", "Hybrid wins:  5 / 12"),
        ("🟡", "Ties:              5 / 12"),
        ("🔴", "Dynamic wins: 2 / 12"),
    ]
    bullet_box(s, score_items, Inches(0.35), Inches(5.75), Inches(4.1), Inches(1.1),
               body_size=14)

    # Main figure: car-s-91 extended dynamics
    img = EXT_PLOTS / "car-s-91_extended.png"
    add_image(s, img, Inches(4.65), Inches(1.28), Inches(5.2))
    add_text_box(s, "car-s-91 — hybrid (orange) vs. dynamic-only (blue) at minimum timeslots",
                 Inches(4.65), Inches(5.72), Inches(5.2), Inches(0.35),
                 font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    # Plateau callout
    callout_box(s,
        "95% of the 10,000-event budget unused — both methods plateau within the first 5%",
        Inches(4.65), Inches(6.1), Inches(8.3), Inches(0.62),
        bg=RGBColor(0xFF, 0xF3, 0xCD), border=GOLD,
        font_size=13, bold=True, color=RGBColor(0x6C, 0x4A, 0x00))

    footer(s)


# ---------------------------------------------------------------------------
# Slide 8 — Sensitivity Analysis
# ---------------------------------------------------------------------------

def slide_sensitivity(prs):
    s = blank_slide(prs)
    slide_bg(s)
    header_bar(s, "Sensitivity Analysis",
               "Small parameter change → large outcome change: deterministic chaos")

    # Left side explanation
    items = [
        ("▸", "Grid of (mode m, drive d) experiments per instance"),
        ("▸", "Measure spread in final clash count → normalized sensitivity score"),
        ("▸", "7 SENSITIVE: small change → large, unpredictable outcome"),
        ("▸", "5 STABLE: result consistent regardless of parameters"),
    ]
    bullet_box(s, items, Inches(0.4), Inches(1.35), Inches(5.5), Inches(2.2),
               body_size=15)

    callout_box(s,
        "Sensitive instances show deterministic chaos:\n"
        "same phenomenon Chirikov (1979) and Strogatz (2001)\n"
        "described in continuous systems — here in a combinatorial space.",
        Inches(0.4), Inches(3.65), Inches(5.5), Inches(1.2),
        bg=RGBColor(0xE8, 0xF0, 0xFE), border=NAVY,
        font_size=13)

    callout_box(s,
        "Practical implication: for the 7 sensitive instances,\n"
        "a short tuning phase would give meaningful improvement.",
        Inches(0.4), Inches(5.0), Inches(5.5), Inches(0.82),
        bg=RGBColor(0xFF, 0xF3, 0xCD), border=GOLD,
        font_size=13)

    # Main figure
    img = SENS / "comparison_normalized.png"
    add_image(s, img, Inches(6.2), Inches(1.28), Inches(6.8))
    add_text_box(s, "Normalized sensitivity score — SENSITIVE vs. STABLE instances",
                 Inches(6.2), Inches(6.45), Inches(6.8), Inches(0.35),
                 font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    footer(s)


# ---------------------------------------------------------------------------
# Slide 9 — Rolling Correlation
# ---------------------------------------------------------------------------

def slide_rolling(prs):
    s = blank_slide(prs)
    slide_bg(s)
    header_bar(s, "Rolling Correlation: Trajectory Divergence",
               "Are the two methods exploring the same space — or genuinely diverging?")

    items = [
        ("▸", "Sliding-window Pearson r between dynamic and hybrid clash histories"),
        ("▸", "Window size W = 20 events; computed over the full 10,000-event run"),
        ("▸", "Initial phase: both methods share a fast shared descent"),
        ("▸", "After Hopfield calls begin: correlation drops toward zero"),
    ]
    bullet_box(s, items, Inches(0.4), Inches(1.35), Inches(5.5), Inches(2.3),
               body_size=15)

    callout_box(s,
        "Mean rolling r ≈ ±0.02 across all 12 instances\n"
        "— statistically indistinguishable from zero",
        Inches(0.4), Inches(3.8), Inches(5.5), Inches(0.85),
        bg=RGBColor(0xD4, 0xED, 0xDA), border=GREEN_OK,
        font_size=15, bold=True, color=RGBColor(0x15, 0x55, 0x24))

    callout_box(s,
        "The two methods explore genuinely different regions of the search space.\n"
        "This justifies calling it a hybrid — not just a modified dynamic walk.",
        Inches(0.4), Inches(4.8), Inches(5.5), Inches(0.85),
        bg=RGBColor(0xE8, 0xF0, 0xFE), border=NAVY,
        font_size=13)

    # Two figures stacked right
    img1 = ROLL / "car-s-91_rolling.png"
    add_image(s, img1, Inches(6.2), Inches(1.28), Inches(6.8))
    add_text_box(s, "car-s-91: clash curves (top) + rolling Pearson r (bottom)",
                 Inches(6.2), Inches(5.25), Inches(6.8), Inches(0.35),
                 font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    # Overview bar chart
    img2 = ROLL / "rolling_correlation_overview.png"
    add_image(s, img2, Inches(6.2), Inches(5.65), Inches(6.8), Inches(1.5))
    add_text_box(s, "Mean rolling r across all 12 instances",
                 Inches(6.2), Inches(7.1), Inches(6.8), Inches(0.3),
                 font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    footer(s)


# ---------------------------------------------------------------------------
# Slide 10 — Conclusions
# ---------------------------------------------------------------------------

def slide_conclusions(prs):
    s = blank_slide(prs)
    slide_bg(s)
    header_bar(s, "Conclusions & Future Work", "")

    # Three takeaway boxes side by side
    box_w = Inches(3.9)
    box_h = Inches(2.2)
    box_top = Inches(1.32)
    gap = Inches(0.2)

    takeaways = [
        ("01", "Physics framework works",
         "Mode selection (quasi-periodic vs. resonant) predicts performance.\n"
         "Not a black box — the dynamics are interpretable."),
        ("02", "Hybrid adds value — where needed",
         "Standard timeslots: dynamic-only sufficient (10/12 solved).\n"
         "Minimum timeslots: Hopfield escape mechanism wins on 5/12."),
        ("03", "Honest limit",
         "At the chromatic number, neither method finds zero clashes in 10,000 events.\n"
         "Restarts or soft constraints are needed."),
    ]

    for i, (num, title, body) in enumerate(takeaways):
        bx = Inches(0.35) + i * (box_w + gap)
        box = s.shapes.add_shape(1, bx, box_top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()

        add_text_box(s, num,
                     bx + Inches(0.12), box_top + Inches(0.1),
                     Inches(0.6), Inches(0.5),
                     font_size=22, bold=True, color=GOLD)
        add_text_box(s, title,
                     bx + Inches(0.12), box_top + Inches(0.52),
                     box_w - Inches(0.24), Inches(0.45),
                     font_size=14, bold=True, color=WHITE)
        add_text_box(s, body,
                     bx + Inches(0.12), box_top + Inches(1.0),
                     box_w - Inches(0.24), Inches(1.1),
                     font_size=11, color=RGBColor(0xCC, 0xD6, 0xE8), wrap=True)

    # Future work section
    add_text_box(s, "Future Work",
                 Inches(0.35), Inches(3.72), Inches(12.5), Inches(0.38),
                 font_size=16, bold=True, color=NAVY)
    add_rect(s, Inches(0.35), Inches(4.12), Inches(12.5), Inches(0.03), GOLD)

    fw_items = [
        "Adaptive Hopfield call frequency (call more when stuck, less when descending)",
        "Restart strategies — multi-start with random perturbation to escape limit cycles",
        "Soft constraint integration → applicable to ITC-2007 and real university deployments",
        "Theoretical analysis of chaos transition as a function of N, S, density, m",
    ]

    fw_top = Inches(4.22)
    for i, item in enumerate(fw_items):
        col = i % 2
        row = i // 2
        add_text_box(s, f"▸  {item}",
                     Inches(0.35) + col * Inches(6.3),
                     fw_top + row * Inches(0.55),
                     Inches(6.1), Inches(0.5),
                     font_size=12, color=DARK_TEXT, wrap=True)

    # Thank you
    callout_box(s,
        "Thank you.  Happy to take questions.",
        Inches(3.0), Inches(6.62), Inches(7.0), Inches(0.6),
        bg=NAVY, border=GOLD, font_size=16, bold=True, color=WHITE)

    footer(s)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = new_prs()

    print("Building slides...")
    slide_title(prs)       ; print("  ✓ Slide 1 — Title")
    slide_problem(prs)     ; print("  ✓ Slide 2 — Problem Statement")
    slide_dataset(prs)     ; print("  ✓ Slide 3 — Toronto Benchmark")
    slide_physics(prs)     ; print("  ✓ Slide 4 — Physics Idea")
    slide_algorithm(prs)   ; print("  ✓ Slide 5 — Algorithm")
    slide_standard(prs)    ; print("  ✓ Slide 6 — Standard Results")
    slide_minimum(prs)     ; print("  ✓ Slide 7 — Minimum Timeslots")
    slide_sensitivity(prs) ; print("  ✓ Slide 8 — Sensitivity Analysis")
    slide_rolling(prs)     ; print("  ✓ Slide 9 — Rolling Correlation")
    slide_conclusions(prs) ; print("  ✓ Slide 10 — Conclusions")

    prs.save(str(OUT_PATH))
    print(f"\nSaved → {OUT_PATH}")


if __name__ == "__main__":
    main()
